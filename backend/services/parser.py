import io
import os
import tempfile
from pypdf import PdfReader
from pptx import Presentation
from fastapi import UploadFile

def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF file."""
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        text = ""
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
        return text
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return ""

def extract_text_from_ppt(file_bytes: bytes) -> str:
    """Extract text from a PowerPoint file."""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error reading PPT: {e}")
        return ""

def save_upload_to_temp(file: UploadFile) -> str:
    """
    保存 UploadFile 到臨時文件

    Args:
        file: FastAPI UploadFile 對象

    Returns:
        臨時文件的絕對路徑
    """
    suffix = os.path.splitext(file.filename)[1]
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(file.file.read())
        return tmp.name

def cleanup_temp_file(path: str) -> None:
    """
    清理臨時文件

    Args:
        path: 文件路徑
    """
    if os.path.exists(path):
        os.remove(path)
